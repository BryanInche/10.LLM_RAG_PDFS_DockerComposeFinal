import os
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
#pip install python-pptx
#pip install python-docx  
#pip install  
from pptx import Presentation
from docx import Document
from langchain.docstore.document import Document as LangchainDocument  # Importar el objeto Document de LangChain
import pypandoc

from atlassian import Jira

#pip install langchain-unstructured


def cargar_documentos(rutas_archivos):
    """
    Carga y divide múltiples documentos en fragmentos de texto.
    Soporta PDFs, PPTs, y Word documents.

    Args:
        rutas_archivos (str o list): Ruta a un archivo, lista de rutas de archivos,
                                     o ruta a una carpeta que contiene archivos.

    Returns:
        list: Lista de documentos divididos en fragmentos.
    """
    todos_los_documentos = []

    # Si se proporciona una sola ruta (str), convertirla en una lista
    if isinstance(rutas_archivos, str):
        rutas_archivos = [rutas_archivos]

    for ruta in rutas_archivos:
        # Si es una carpeta, obtener todos los archivos compatibles dentro de ella
        if os.path.isdir(ruta):
            print(f"Procesando carpeta: {ruta}")
            archivos_en_carpeta = [
                os.path.join(ruta, archivo) for archivo in os.listdir(ruta)
                if archivo.endswith(".pdf") or archivo.endswith(".ppt") or archivo.endswith(".docx")
            ]
            if not archivos_en_carpeta:
                print(f"Advertencia: No se encontraron archivos compatibles en la carpeta {ruta}.")
            else:
                print(f"Se encontraron {len(archivos_en_carpeta)} archivos compatibles en la carpeta {ruta}.")
            rutas_archivos.extend(archivos_en_carpeta)
            continue

        # Si es un archivo PDF, cargarlo
        if ruta.endswith(".pdf"):
            try:
                print(f"Procesando archivo PDF: {ruta}")
                carga = PyPDFLoader(ruta)
                #carga = UnstructuredFileLoader(ruta) # Para PDFs problemáticos
                documentos = carga.load()
            except Exception as e:
                print(f"Error al cargar el archivo PDF {ruta}: {e}")
                continue

        # Si es un archivo PPTX, cargarlo
        elif ruta.endswith(".ppt"):
            try:
                print(f"Procesando archivo PPT: {ruta}")
                ############################################
                #prs = Presentation(ruta)
                #textos = []
                #for slide in prs.slides:
                #    for shape in slide.shapes:
                #        if hasattr(shape, "text"):
                #            textos.append(shape.text)
                ##############################################
                texto = pypandoc.convert_file(ruta, 'plain') #extraer el texto directamente desde los archivos .ppt

                # Crear un objeto Document de LangChain
                documentos = [LangchainDocument(page_content=texto)]            
                # Crear un objeto Document de LangChain
                #documentos = [LangchainDocument(page_content="\n".join(textos))]
            except Exception as e:
                print(f"Error al cargar el archivo PPT {ruta}: {e}")
                continue

        # Si es un archivo DOCX, cargarlo
        elif ruta.endswith(".docx"):
            try:
                print(f"Procesando archivo DOCX: {ruta}")
                doc = Document(ruta)
                textos = [para.text for para in doc.paragraphs]
                # Crear un objeto Document de LangChain
                documentos = [LangchainDocument(page_content="\n".join(textos))]
            except Exception as e:
                print(f"Error al cargar el archivo DOCX {ruta}: {e}")
                continue

        else:
            print(f"Advertencia: La ruta {ruta} no es un archivo compatible.")
            continue

        if not documentos:
            print(f"Advertencia: No se pudieron cargar documentos desde el archivo {ruta}.")
            continue

        # Dividir los documentos en fragmentos
        try:
            texto_spliter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=800)
            documentos_divididos = texto_spliter.split_documents(documentos)
        except Exception as e:
            print(f"Error al dividir los documentos del archivo {ruta}: {e}")
            continue

        # Agregar los documentos divididos a la lista general
        todos_los_documentos.extend(documentos_divididos)
        print(f"Se cargaron {len(documentos_divididos)} fragmentos desde {ruta}.")

    if not todos_los_documentos:
        raise ValueError("No se pudieron cargar documentos desde ninguno de los archivos o carpetas proporcionados.")

    print(f"Se cargaron un total de {len(todos_los_documentos)} fragmentos de texto.")
    return todos_los_documentos


def crear_vectorstore(docs):
    """
    Crea un vector store a partir de una lista de documentos.

    Args:
        docs (list): Lista de documentos divididos en fragmentos.

    Returns:
        Chroma: Vector store creado.
    """
    # Validar que se proporcionen documentos
    if not docs:
        raise ValueError("No se proporcionaron documentos para crear el vector store.")

    # Crear el modelo de embeddings
    try:
        # Busqueda semantica, basada en similitud de embedings
        # Modelo compatible con FastEmbedEmbeddings
        #embeding_modelo = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")
        
        # Modelo de Hugging Face
        embeding_modelo = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    except Exception as e:
        print(f"Error al crear el modelo de embeddings: {e}")
        raise

    # Crear el vector store
    try:
        vector_store = Chroma.from_documents(
            documents=docs,
            embedding=embeding_modelo,
            persist_directory="4chroma_ms4m_bd_allche",
            collection_name="4data_ms4m_allche" #colection dentro del vs
        )
        print("Vector store creado correctamente.")
    except Exception as e:
        print(f"Error al crear el vector store: {e}")
        raise

    # Depuración: Verificar el número de documentos en el vector store
    print(f"Número de documentos en el vector store: {vector_store._collection.count()}")

    # Validar que el vector store no esté vacío
    if vector_store._collection.count() == 0:
        raise ValueError("El vector store está vacío. No se cargaron documentos.")

    # Realizar una búsqueda de prueba
    try:
        query = "Area de planeamiento de ms4m"  # Cambia esto por una consulta relevante para tu caso
        resultados = vector_store.similarity_search(query, k=3)
        print(f"Búsqueda de prueba realizada. Resultados encontrados: {len(resultados)}")
    except Exception as e:
        print(f"Error al realizar la búsqueda de prueba: {e}")
        raise

    return vector_store



# Creamos el VectorStore con los nuevos Documentos, Importante que este Vector Store se cree cada vez que deseemos agregar nuevos documentos
#ruta_files = [
#    #"pdf_solos/ASIS_REPORTE_EJECUTIVO.pdf",  # Ruta a un archivo PDF
#    "ppt_solos",  # Ruta a una carpeta de PPTs
#    "word_solos",  # Ruta a una carpeta de Words
#    "pdf_carpetas/pdfs_varios"  # Ruta a una carpeta de PDFs
#]

### Cargar documentos
#documentos = cargar_documentos(ruta_files)

# Crear el vector store
#vector_store = crear_vectorstore(documentos)

